import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Inches
from pptx.slide import Slide
from PIL import Image, UnidentifiedImageError

import logging
from typing import Literal, Union, List, Dict, Any

ChartTypes = Literal["bar", "line", "pie", "scatter", "area"]

class PresentationManager:
    # Slide layout constants
    SLIDE_LAYOUT_TITLE = 0
    SLIDE_LAYOUT_TITLE_AND_CONTENT = 1
    SLIDE_LAYOUT_SECTION_HEADER = 2
    SLIDE_LAYOUT_TWO_CONTENT = 3
    SLIDE_LAYOUT_COMPARISON = 4
    SLIDE_LAYOUT_TITLE_ONLY = 5
    SLIDE_LAYOUT_BLANK = 6
    SLIDE_LAYOUT_CONTENT_WITH_CAPTION = 7
    SLIDE_LAYOUT_PICTURE_WITH_CAPTION = 8


    def __init__(self):
        self.presentations: Dict[str, Any] = {}

    def _add_formatted_bullets(self, text_frame, text_block):
        """
        Process a text block and add paragraphs with proper bullet indentation
        using ASCII code detection:
        - ASCII 10 (LF) or ASCII 13 (CR) or combination for new lines (main bullets)
        - ASCII 9 (HT) for tab indentation (sub-bullets)

        Args:
            text_frame: The PowerPoint text frame to add text to
            text_block: String of text to process
        """
        # First, normalize all line endings to a single format
        # Replace CR+LF (Windows) with a single marker
        normalized_text = text_block.replace('\r\n', '\n')
        # Replace any remaining CR (old Mac) with LF
        normalized_text = normalized_text.replace('\r', '\n')

        # Split the text block into lines using ASCII 10 (LF)
        lines = normalized_text.split('\n')

        # Clear any existing text
        if text_frame.paragraphs:
            p = text_frame.paragraphs[0]
            p.text = ""
        else:
            p = text_frame.add_paragraph()

        # Process the first line separately (if it exists)
        if lines and lines[0].strip():
            first_line = lines[0]
            # Count leading tabs (ASCII 9) to determine indentation level
            level = 0
            while first_line and ord(first_line[0]) == 9:  # ASCII 9 is HT (tab)
                level += 1
                first_line = first_line[1:]

            p.text = first_line.strip()
            p.level = level

        # Process remaining lines
        for line in lines[1:]:
            if not line.strip():
                continue  # Skip empty lines

            # Count leading tabs (ASCII 9) to determine indentation level
            level = 0
            while line and ord(line[0]) == 9:  # ASCII 9 is HT (tab)
                level += 1
                line = line[1:]

            # Add the paragraph with proper indentation
            p = text_frame.add_paragraph()
            p.text = line.strip()
            p.level = level

    def add_section_header_slide(self, presentation_name: str, header: str, subtitle: str):
        """
        Create a section header slide for the given presentation

        Args:
            presentation_name: The presentation to add the slide to
            header: The section header to use
            subtitle: The subtitle of the section header to use
        """
        try:
            prs = self.presentations[presentation_name]
        except KeyError as e:
            raise ValueError(f"Presentation '{presentation_name}' not found")
        slide_master = prs.slide_master

        # Add a new slide with layout
        slide_layout = prs.slide_layouts[self.SLIDE_LAYOUT_SECTION_HEADER]
        slide = prs.slides.add_slide(slide_layout)

        # Set the subtitle
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            text_frame = subtitle_shape.text_frame
            text_frame.text = subtitle

        # Set the section header
        if header:
            header_shape = slide.shapes.title
            header_shape.text = header

        return slide

    def add_comparison_slide(self, presentation_name: str, title: str, left_side_title: str, left_side_content: str,
                             right_side_title: str, right_side_content: str ):
        """
        Create a section header slide for the given presentation

        Args:
            presentation_name: The presentation to add the slide to
            title: The title of the slide
            left_side_title: The title of the left hand side content
            left_side_content: The body content for the left hand side
            right_side_title: The title of the right hand side content
            right_side_content: The body content for the right hand side
        """
        try:
            prs = self.presentations[presentation_name]
        except KeyError as e:
            raise ValueError(f"Presentation '{presentation_name}' not found")
        slide_master = prs.slide_master

        # Add a new slide with layout
        slide_layout = prs.slide_layouts[self.SLIDE_LAYOUT_COMPARISON]
        slide = prs.slides.add_slide(slide_layout)

        # Set the title
        title_shape = slide.shapes.title
        title_shape.text = title

        # Build the left hand content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.text = left_side_title

        content_shape = slide.placeholders[2]
        text_frame = content_shape.text_frame
        text_frame.text = left_side_content

        # Build the right hand content
        content_shape = slide.placeholders[3]
        text_frame = content_shape.text_frame
        text_frame.text = right_side_title

        content_shape = slide.placeholders[4]
        text_frame = content_shape.text_frame
        text_frame.text = right_side_content
        return slide

    def add_picture_with_caption_slide(self, presentation_name: str, title: str,
                                       image_path: str, caption_text: str) -> Slide:

        """
        For the given presentation builds a slide with the picture with caption template.
        Maintains the image's aspect ratio by adjusting the picture object after insertion.
        Args:
            presentation_name: The presentation to add the slide to
            title: The title of the slide
            image_path: The path to the image to insert
            caption_text: The caption content

        """
        try:
            prs = self.presentations[presentation_name]
        except KeyError as e:
            raise ValueError(f"Presentation '{presentation_name}' not found")

        # Add a new slide with layout 8 (Picture with Caption)
        try:
            slide_layout = prs.slide_layouts[self.SLIDE_LAYOUT_PICTURE_WITH_CAPTION]
            slide = prs.slides.add_slide(slide_layout)
        except IndexError as e:
            error_message = f"Slide Index does not exist. Error: {str(e)}"
            raise ValueError(error_message)
        # Set the title
        title_shape = slide.shapes.title
        title_shape.text = title

        # Get the image placeholder
        try:
            placeholder = slide.placeholders[1]
        except IndexError as e:
            error_message = f"Placeholder index does not exist. Error {str(e)}"
            raise ValueError(error_message)

        # Insert the picture into the placeholder
        if not os.path.exists(image_path):
            raise FileNotFoundError(f"Image not found: {image_path}")
        try:
            picture = placeholder.insert_picture(image_path)
        except FileNotFoundError as e:
            error_message = f"Image not found during insertion: {str(e)}"
            raise
        except UnidentifiedImageError as e:
            error_message = f"Image file {image_path} is not a valid image: {str(e)}"
            raise ValueError(error_message)
        except Exception as e:
            error_message = f"An unexpected error occured during picture insertion: {str(e)}"
            raise

        # Get placeholder dimensions after picture insertion
        available_width = picture.width
        available_height = picture.height

        # Get original image dimensions directly from the picture object
        image_width, image_height = picture.image.size

        # Calculate aspect ratios
        placeholder_aspect_ratio = float(available_width) / float(available_height)
        image_aspect_ratio = float(image_width) / float(image_height)

        # Store initial position
        pos_left, pos_top = picture.left, picture.top

        # Remove any cropping
        picture.crop_top = 0
        picture.crop_left = 0
        picture.crop_bottom = 0
        picture.crop_right = 0

        # Adjust picture dimensions based on aspect ratio comparison
        if placeholder_aspect_ratio > image_aspect_ratio:
            # Placeholder is wider than image - adjust width down while maintaining height
            picture.width = int(image_aspect_ratio * available_height)
            picture.height = available_height
        else:
            # Placeholder is taller than image - adjust height down while maintaining width
            picture.height = int(available_width / image_aspect_ratio)
            picture.width = available_width

        # Center the image within the available space
        picture.left = pos_left + int((available_width - picture.width) / 2)
        picture.top = pos_top + int((available_height - picture.height) / 2)

        # Set the caption
        caption = slide.placeholders[2]
        caption.text = caption_text

        return slide

    def add_title_with_content_slide(self, presentation_name: str, title: str, content: str) -> Slide:
        try:
            prs = self.presentations[presentation_name]
        except KeyError as e:
            raise ValueError(f"Presentation '{presentation_name}' not found")
        slide_master = prs.slide_master
        # Add a slide with title and content
        slide_layout = prs.slide_layouts[self.SLIDE_LAYOUT_TITLE_AND_CONTENT]  # Use layout with title and content
        slide = prs.slides.add_slide(slide_layout)

        # Set the title
        title_shape = slide.shapes.title
        title_shape.text = title

        # Set the content
        content_shape = slide.placeholders[1]
        #content_shape.text = content
        # Get the content placeholder and add our formatted text

        text_frame = content_shape.text_frame
        self._add_formatted_bullets(text_frame, content)
        return slide

    def add_table_slide(self, presentation_name: str, title: str, headers: str, rows: str) -> Slide:

        try:
            prs = self.presentations[presentation_name]
        except KeyError as e:
            raise ValueError(f"Presentation '{presentation_name}' not found")

        slide_layout = prs.slide_layouts[self.SLIDE_LAYOUT_TITLE_ONLY]
        slide = prs.slides.add_slide(slide_layout)

        # Set the title
        title_shape = slide.shapes.title
        title_shape.text = title

        # Calculate table dimensions and position
        num_rows = len(rows) + 1  # +1 for header row
        num_cols = len(headers)

        # Position table in the middle of the slide with some margins
        x = Inches(1)  # Left margin
        y = Inches(2)  # Top margin below title

        # Make table width proportional to the number of columns
        width_per_col = Inches(8 / num_cols)  # Divide available width (8 inches) by number of columns
        height_per_row = Inches(0.4)  # Standard height per row

        # Create table
        shape = slide.shapes.add_table(
            num_rows,
            num_cols,
            x,
            y,
            width_per_col * num_cols,
            height_per_row * num_rows
        )
        table = shape.table

        # Add headers
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            # Style header row
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = Pt(11)



        # Add data rows
        for row_idx, row_data in enumerate(rows, start=1):
            for col_idx, cell_value in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_value)
                # Style data cells
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(10)

        return slide

    def add_title_slide(self, presentation_name: str, title: str) -> Slide:
        try:
            prs = self.presentations[presentation_name]
        except KeyError as e:
            raise ValueError(f"Presentation '{presentation_name}' not found")

        # Add a slide with title and content
        slide_layout = prs.slide_layouts[self.SLIDE_LAYOUT_TITLE]
        slide = prs.slides.add_slide(slide_layout)

        # Set the title
        title_shape = slide.shapes.title
        title_shape.text = title
        return slide


