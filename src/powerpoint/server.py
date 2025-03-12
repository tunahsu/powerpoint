import os
from mcp.server import Server, NotificationOptions
from mcp.server.models import InitializationOptions
import mcp.server.stdio
import mcp.types as types
import asyncio
from pptx import Presentation
import logging
from .presentation_manager import PresentationManager
from .chart_manager import ChartManager
from .vision_manager import VisionManager

logger = logging.getLogger('mcp_powerpoint_server')
logger.info("Starting MCP Powerpoint Server")

BACKUP_FILE_NAME = 'backup.pptx'

def sanitize_path(base_path: str, file_name: str) -> str:
    """
    Ensure that the resulting path doesn't escape outside the base directory
    Returns a safe, normalized path
    """

    joined_path = os.path.join(base_path, file_name)
    normalized_path = os.path.normpath(joined_path)

    if not normalized_path.startswith(base_path):
        raise ValueError(f"Invalid path. Attempted to access location outside allowed directory.")

    return normalized_path

async def main(folder_path):
    logger.info(f"Starting Powerpoint MCP Server")
    presentation_manager = PresentationManager()
    chart_manager = ChartManager()
    vision_manager = VisionManager()
    server = Server("powerpoint-server")
    logger.debug("Registering Handlers")
    path = folder_path


    @server.list_tools()
    async def handle_list_tools() -> list[types.Tool]:
        """List available PowerPoint tools."""
        return [
            types.Tool(
                name="create-presentation",
                description="This tool starts the process of generating a new powerpoint presentation with the name given "
                            "by the user. Use this tool when the user requests to create or generate a new presentation.",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "name": {
                            "type": "string",
                            "description": "Name of the presentation (without .pptx extension)",
                        },
                    },
                    "required": ["name"],
                },
            ),
            types.Tool(
                name="generate-and-save-image",
                description="Generates an image using a FLUX model and save the image to the specified path. The tool "
                            "will return a PNG file path. It should be used when the user asks to generate or create an "
                            "image or a picture.",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "prompt": {
                            "type": "string",
                            "description": "Description of the image to generate in the form of a prompt.",
                        },
                        "file_name": {
                            "type": "string",
                            "description": "Filename of the image. Include the extension of .png",
                        },
                    },
                    "required": ["prompt", "file_name"],
                },
            ),
            types.Tool(
                name="add-slide-title-only",
                description="This tool adds a new title slide to the presentation you are working on. The tool doesn't "
                            "return anything. It requires the presentation_name to work on.",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "presentation_name": {
                            "type": "string",
                            "description": "Name of the presentation to add the slide to",
                        },
                        "title": {
                            "type": "string",
                            "description": "Title of the slide",
                        }
                    },
                    "required": ["presentation_name", "title"],
                },
            ),
            types.Tool(
                name="add-slide-section-header",
                description="This tool adds a section header (a.k.a segue) slide to the presentation you are working on. The tool doesn't "
                            "return anything. It requires the presentation_name to work on.",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "presentation_name": {
                            "type": "string",
                            "description": "Name of the presentation to add the slide to",
                        },
                        "header": {
                            "type": "string",
                            "description": "Section header title",
                        },
                        "subtitle": {
                            "type": "string",
                            "description": "Section header subtitle",
                        }

                    },
                    "required": ["presentation_name", "header"],
                },
            ),
            types.Tool(
                name="add-slide-title-content",
                description="Add a new slide with a title and content to an existing presentation",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "presentation_name": {
                            "type": "string",
                            "description": "Name of the presentation to add the slide to",
                        },
                        "title": {
                            "type": "string",
                            "description": "Title of the slide",
                        },
                        "content": {
                            "type": "string",
                            "description": "Content/body text of the slide. "
                                           "Separate main points with a single carriage return character."
                                           "Make sub-points with tab character."
                                           "Do not use bullet points, asterisks or dashes for points."
                                           "Max main points is 4"
                        },
                    },
                    "required": ["presentation_name", "title", "content"],
                },
            ),
            types.Tool(
                name="add-slide-comparison",
                description="Add a new a comparison slide with title and comparison content. Use when you wish to "
                            "compare two concepts",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "presentation_name": {
                            "type": "string",
                            "description": "Name of the presentation to add the slide to",
                        },
                        "title": {
                            "type": "string",
                            "description": "Title of the slide",
                        },
                        "left_side_title": {
                            "type": "string",
                            "description": "Title of the left concept",
                        },
                        "left_side_content": {
                            "type": "string",
                            "description": "Content/body text of left concept. "
                                           "Separate main points with a single carriage return character."
                                           "Make sub-points with tab character."
                                           "Do not use bullet points, asterisks or dashes for points."
                                           "Max main points is 4"
                        },
                        "right_side_title": {
                            "type": "string",
                            "description": "Title of the right concept",
                        },
                        "right_side_content": {
                            "type": "string",
                            "description": "Content/body text of right concept. "
                                           "Separate main points with a single carriage return character."
                                           "Make sub-points with tab character."
                                           "Do not use bullet points, asterisks or dashes for points."
                                           "Max main points is 4"
                        },
                    },
                    "required": ["presentation_name", "title", "left_side_title", "left_side_content",
                                 "right_side_title", "right_side_content"],
                },
            ),
            types.Tool(
                name="add-slide-title-with-table",
                description="Add a new slide with a title and table containing the provided data",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "presentation_name": {
                            "type": "string",
                            "description": "Name of the presentation to add the slide to",
                        },
                        "title": {
                            "type": "string",
                            "description": "Title of the slide",
                        },
                        "data": {
                            "type": "object",
                            "description": "Table data object with headers and rows",
                            "properties": {
                                "headers": {
                                    "type": "array",
                                    "items": {"type": "string"},
                                    "description": "Array of column headers"
                                },
                                "rows": {
                                    "type": "array",
                                    "items": {
                                        "type": "array",
                                        "items": {"type": ["string", "number"]},
                                    },
                                    "description": "Array of row data arrays"
                                }
                            },
                            "required": ["headers", "rows"]
                        }
                    },
                    "required": ["presentation_name", "title", "data"],
                },
            ),
            types.Tool(
                name="add-slide-title-with-chart",
                description="Add a new slide with a title and chart. The chart type will be automatically selected based on the data structure.",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "presentation_name": {
                            "type": "string",
                            "description": "Name of the presentation to add the slide to",
                        },
                        "title": {
                            "type": "string",
                            "description": "Title of the slide",
                        },
                        "data": {
                            "type": "object",
                            "description": "Chart data structure",
                            "properties": {
                                "categories": {
                                    "type": "array",
                                    "items": {"type": ["string", "number"]},
                                    "description": "X-axis categories or labels (optional)"
                                },
                                "series": {
                                    "type": "array",
                                    "items": {
                                        "type": "object",
                                        "properties": {
                                            "name": {
                                                "type": "string",
                                                "description": "Name of the data series"
                                            },
                                            "values": {
                                                "type": "array",
                                                "items": {
                                                    "oneOf": [
                                                        {"type": "number"},
                                                        {
                                                            "type": "array",
                                                            "items": {"type": "number"},
                                                            "minItems": 2,
                                                            "maxItems": 2
                                                        }
                                                    ]
                                                },
                                                "description": "Values for the series. Can be simple numbers or [x,y] pairs for scatter plots"
                                            }
                                        },
                                        "required": ["name", "values"]
                                    }
                                },
                                "x_axis": {
                                    "type": "string",
                                    "description": "X-axis title (optional)"
                                },
                                "y_axis": {
                                    "type": "string",
                                    "description": "Y-axis title (optional)"
                                }
                            },
                            "required": ["series"]
                        }
                    },
                    "required": ["presentation_name", "title", "data"],
                },
            ),
            types.Tool(
                name="add-slide-picture-with-caption",
                description="Add a new slide with a picture and caption to an existing presentation",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "presentation_name": {
                            "type": "string",
                            "description": "Name of the presentation to add the slide to",
                        },
                        "title": {
                            "type": "string",
                            "description": "Title of the slide",
                        },
                        "caption": {
                            "type": "string",
                            "description": "Caption text to appear below the picture"
                        },
                        "image_path": {
                            "type": "string",
                            "description": "Path to the image file to insert"
                        }
                    },
                    "required": ["presentation_name", "title", "caption", "image_path"],
                },
            ),
            types.Tool(
                name="open-presentation",
                description="Opens an existing presentation and saves a copy to a new file for backup. Use this tool when "
                            "the user requests to open a presentation that has already been created.",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "presentation_name": {
                            "type": "string",
                            "description": "Name of the presentation to open",
                        },
                        "output_path": {
                            "type": "string",
                            "description": "Path where to save the presentation (optional)",
                        },
                    },
                    "required": ["presentation_name"],
                },
            ),
            types.Tool(
                name="save-presentation",
                description="Save the presentation to a file. Always use this tool at the end of any process that has "
                            "added slides to a presentation.",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "presentation_name": {
                            "type": "string",
                            "description": "Name of the presentation to save",
                        },
                        "output_path": {
                            "type": "string",
                            "description": "Path where to save the presentation (optional)",
                        },
                    },
                    "required": ["presentation_name"],
                },
            ),
        ]


    @server.call_tool()
    async def handle_call_tool(
            name: str, arguments: dict | None
    ) -> list[types.TextContent | types.ImageContent | types.EmbeddedResource]:
        """Handle PowerPoint tool execution requests."""
        if not arguments:
            raise ValueError("Missing arguments")
        if name == "open-presentation":
            presentation_name = arguments.get("presentation_name")
            if not presentation_name:
                raise ValueError("Missing presentation name")
            file_name = f"{presentation_name}.pptx"

            try:
                safe_file_path = sanitize_path(folder_path, file_name)
            except ValueError as e:
                raise ValueError(f"Invalid file path: {str(e)}")

            # attempt to load presentation
            try:
                prs = Presentation(safe_file_path)
            except Exception as e:
                raise ValueError(f"Unable to load {safe_file_path}. Error: {str(e)}")

            # Create a backup of the original file
            file_name = BACKUP_FILE_NAME
            try:
                safe_file_path = sanitize_path(folder_path, file_name)
            except ValueError as e:
                raise ValueError(f"Invalid file path: {str(e)}")
            # attempt to save a backup of presentation
            try:
                prs.save(safe_file_path)
            except Exception as e:
                raise ValueError(f"Unable to save {safe_file_path}. Error: {str(e)}")

            presentation_manager.presentations[presentation_name] = prs

            return [
                types.TextContent(
                    type="text",
                    text=f"Opened presentation: {presentation_name}"
                )
            ]
        elif name == "generate-and-save-image":
            prompt = arguments.get("prompt")
            file_name = arguments.get("file_name")
            try:
                safe_file_path = sanitize_path(folder_path, file_name)
            except ValueError as e:
                raise ValueError(f"Invalid file path: {str(e)}")

            if not all([prompt, file_name]):
                raise ValueError("Missing required arguments")

            try:
                saved_path = await vision_manager.generate_and_save_image(prompt, str(safe_file_path))
                return [
                    types.TextContent(
                        type="text",
                        text=f"Successfully generated and saved image to: {saved_path}"
                    )
                ]
            except Exception as e:
                return [
                    types.TextContent(
                        type="text",
                        text=f"Failed to generate image: {str(e)}"
                    )
                ]
        elif name == "add-slide-comparison":
            # Get arguments
            presentation_name = arguments["presentation_name"]
            title = arguments["title"]
            left_side_title = arguments["left_side_title"]
            left_side_content = arguments["left_side_content"]
            right_side_title = arguments["right_side_title"]
            right_side_content = arguments["right_side_content"]

            if not all([presentation_name, title, left_side_title, left_side_content,
                        right_side_title, right_side_content]):
                raise ValueError("Missing required arguments")

            if presentation_name not in presentation_manager.presentations:
                raise ValueError(f"Presentation not found: {presentation_name}")
            try:
                slide = presentation_manager.add_comparison_slide(presentation_name, title, left_side_title,
                                                                  left_side_content, right_side_title, right_side_content)
            except Exception as e:
                raise ValueError(f"Unable to add comparison slide to {presentation_name}.pptx")

            return [types.TextContent(
                type="text",
                text=f"Successfully added comparison slide {title} to {presentation_name}.pptx"
            )]

        elif name == "add-slide-picture-with-caption":

            # Get arguments
            presentation_name = arguments["presentation_name"]
            title = arguments["title"]
            caption = arguments["caption"]
            file_name = arguments["image_path"]

            if not all([presentation_name, title, caption, file_name]):
                raise ValueError("Missing required arguments")

            if presentation_name not in presentation_manager.presentations:
                raise ValueError(f"Presentation not found: {presentation_name}")

            try:
                safe_file_path = sanitize_path(folder_path, file_name)
            except ValueError as e:
                raise ValueError(f"Invalid file path: {str(e)}")

            try:
                slide = presentation_manager.add_picture_with_caption_slide(presentation_name, title, str(safe_file_path), caption)
            except Exception as e:
                raise ValueError(f"Unable to add slide with caption and picture layout to {presentation_name}.pptx. Error: {str(e)}")

            return [types.TextContent(
                type="text",
                text=f"Successfully added slide with caption and picture layout to {presentation_name}.pptx"
            )]

        elif name == "create-presentation":

            presentation_name = arguments.get("name")
            if not presentation_name:
                raise ValueError("Missing presentation name")

            # Create new presentation
            prs = Presentation()
            try:
                presentation_manager.presentations[presentation_name] = prs
            except KeyError as e:
                raise ValueError(f"Unable to add {presentation_name} to presentation. Error: {str(e)}")

            return [
                types.TextContent(
                    type="text",
                    text=f"Created new presentation: {presentation_name}"
                )
            ]

        elif name == "add-slide-title-content":
            presentation_name = arguments.get("presentation_name")
            title = arguments.get("title")
            content = arguments.get("content")

            if not all([presentation_name, title, content]):
                raise ValueError("Missing required arguments")

            if presentation_name not in presentation_manager.presentations:
                raise ValueError(f"Presentation not found: {presentation_name}")

            try:
                slide = presentation_manager.add_title_with_content_slide(presentation_name, title, content)
            except Exception as e:
                raise ValueError(f"Unable to add slide '{title}' to presentation: {presentation_name}")

            return [
                types.TextContent(
                    type="text",
                    text=f"Added slide '{title}' to presentation: {presentation_name}"
                )
            ]
        elif name == "add-slide-section-header":
            presentation_name = arguments.get("presentation_name")
            header = arguments.get("header")
            subtitle = arguments.get("subtitle")

            if not all([presentation_name, header]):
                raise ValueError("Missing required arguments")

            if presentation_name not in presentation_manager.presentations:
                raise ValueError(f"Presentation not found: {presentation_name}")

            try:
                slide = presentation_manager.add_section_header_slide(presentation_name, header, subtitle)
            except Exception as e:
                raise ValueError(f"Unable to add slide '{header}' to presentation: {presentation_name}")

            return [
                types.TextContent(
                    type="text",
                    text=f"Added slide '{header}' to presentation: {presentation_name}"
                )
            ]
        elif name == "add-slide-title-with-table":
            presentation_name = arguments.get("presentation_name")
            title = arguments.get("title")
            table_data = arguments.get("data")

            if not all([presentation_name, title, table_data]):
                raise ValueError("Missing required arguments")

            if presentation_name not in presentation_manager.presentations:
                raise ValueError(f"Presentation not found: {presentation_name}")

            # Validate table data structure
            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])

            if not headers:
                raise ValueError("Table headers are required")

            if not rows:
                raise ValueError("Table rows are required")

            # Validate that all rows match header length
            if not all(len(row) == len(headers) for row in rows):
                raise ValueError("All rows must have the same number of columns as headers")
            try:
                slide = presentation_manager.add_table_slide(presentation_name, title, headers, rows)
            except Exception as e:
                raise ValueError(f"Unable to add slide '{title}' with a table to presentation: {presentation_name}")

            return [
                types.TextContent(
                    type="text",
                    text=f"Added slide '{title}' with a table to presentation: {presentation_name}"
                )
            ]
        elif name == "add-slide-title-with-chart":
            presentation_name = arguments.get("presentation_name")
            title = arguments.get("title")
            chart_data = arguments.get("data")

            if not all([presentation_name, title, chart_data]):
                raise ValueError("Missing required arguments")

            if presentation_name not in presentation_manager.presentations:
                raise ValueError(f"Presentation not found: {presentation_name}")

            # Get the presentation and create a new slide
            prs = presentation_manager.presentations[presentation_name]
            slide_layout = prs.slide_layouts[5]  # Title and blank content
            slide = prs.slides.add_slide(slide_layout)

            # Set the title
            title_shape = slide.shapes.title
            title_shape.text = title

            # Determine the best chart type for the data
            try:
                chart_type, chart_format = chart_manager.determine_chart_type(chart_data)
            except Exception as e:
                raise ValueError(f"Unable to determine chart type.")

            # Add the chart to the slide
            try:
                chart = chart_manager.add_chart_to_slide(slide, chart_type, chart_data, chart_format)
                chart_type_name = chart_type.name.lower().replace('xl_chart_type.', '')

                return [
                    types.TextContent(
                        type="text",
                        text=f"Added slide '{title}' with a {chart_type_name} chart to presentation: {presentation_name}"
                    )
                ]
            except Exception as e:
                raise ValueError(f"Failed to create slide with chart: {str(e)}")
        elif name == "add-slide-title-only":
            presentation_name = arguments.get("presentation_name")
            title = arguments.get("title")

            if not all([presentation_name, title]):
                raise ValueError("Missing required arguments")

            if presentation_name not in presentation_manager.presentations:
                raise ValueError(f"Presentation not found: {presentation_name}")

            try:
                slide = presentation_manager.add_title_slide(presentation_name, title)
            except Exception as e:
                 raise ValueError(f"Unable to add '{title} to presentation: {presentation_name}. Error: {e}")

            return [
                types.TextContent(
                    type="text",
                    text=f"Added slide '{title}' to presentation: {presentation_name}"
                )
            ]
        elif name == "save-presentation":
            presentation_name = arguments.get("presentation_name")
            output_path = arguments.get("output_path")


            if not presentation_name:
                raise ValueError("Missing presentation name")

            if presentation_name not in presentation_manager.presentations:
                raise ValueError(f"Presentation not found: {presentation_name}")

            prs = presentation_manager.presentations[presentation_name]

            # Default output path if none provided
            if not output_path:
                output_path = f"{presentation_name}.pptx"

            file_path = os.path.join(path,output_path)
            # Save the presentation
            try:
                prs.save(file_path)
            except Exception as e:
                raise ValueError(f"Unable to save the {presentation_name}. Error: {e}")

            return [
                types.TextContent(
                    type="text",
                    text=f"Saved presentation to: {file_path}"
                )
            ]

        else:
            raise ValueError(f"Unknown tool: {name}")

    async with mcp.server.stdio.stdio_server() as (read_stream, write_stream):
        logger.info("Server running with stdio transport")
        await server.run(
            read_stream,
            write_stream,
            InitializationOptions(
                server_name="powerpoint",
                server_version="0.1.0",
                capabilities=server.get_capabilities(
                    notification_options=NotificationOptions(),
                    experimental_capabilities={},
                ),
            ),
        )


if __name__ == "__main__":
    asyncio.run(main())