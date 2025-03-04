from pptx.chart import chart
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_LEGEND_POSITION
from typing import Literal, Union, List, Dict, Any

class ChartManager:
    def __init__(self):
        self.name = "Chart Manager"

    def determine_chart_type(self, data: Dict[str, Any]) -> tuple[XL_CHART_TYPE, str]:
        """
        Analyze the data structure and determine the most appropriate chart type.
        Returns tuple of (PowerPoint chart type enum, chart_format)
        """
        # evaluate the data
        series_count = len(data["series"])
        categories = data.get("categories", [])

        # Check for XY data more safely by checking the first value of each series
        is_xy_data = False
        for series in data["series"]:
            values = series.get("values", [])
            if values:
                first_value = values[0]
                is_xy_data = isinstance(first_value, (list, tuple)) and len(first_value) == 2
                break

        if is_xy_data:
            return XL_CHART_TYPE.XY_SCATTER, "xy"

        # If we have percentage data that adds up to ~100, suggest pie chart
        if series_count == 1 and categories:
            values = data["series"][0].get("values", [])
            if len(values) <= 8:
                try:
                    total = sum(float(v) for v in values)
                    if 95 <= total <= 105:
                        return XL_CHART_TYPE.PIE, "category"
                except (TypeError, ValueError):
                    pass

        # For time series or trending data, suggest line chart
        if categories and any(
                isinstance(cat, (str, int)) and
                any(term in str(cat).lower() for term in
                    ["date", "time", "year", "month", "quarter", "q1", "q2", "q3", "q4"])
                for cat in categories
        ):
            return XL_CHART_TYPE.LINE, "category"

        # For multiple series comparing values, suggest bar chart
        if series_count > 1 and categories:
            return XL_CHART_TYPE.BAR_CLUSTERED, "category"

        # Default to column chart for single series
        return XL_CHART_TYPE.COLUMN_CLUSTERED, "category"


    def add_chart_to_slide(self, slide, chart_type: XL_CHART_TYPE, data: Dict[str, Any],
                           chart_format: str = "category") -> chart:
        """Add a chart to the slide with the specified data."""
        # Position chart in the middle of the slide with margins
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(5)

        if chart_format == "category":
            chart_data = CategoryChartData()
            chart_data.categories = data.get("categories", [])

            # Add each series
            for series in data["series"]:
                chart_data.add_series(series["name"], series["values"])

        elif chart_format == "xy":
            chart_data = XyChartData()

            # Add each series
            for series in data["series"]:
                series_data = chart_data.add_series(series["name"])
                for x, y in series["values"]:
                    series_data.add_data_point(x, y)

        # Add and configure the chart
        graphic_frame = slide.shapes.add_chart(
            chart_type, left, top, width, height, chart_data
        )
        chart = graphic_frame.chart

        # Basic formatting
        chart.has_legend = True
        if len(data["series"]) > 1:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        # Add axis titles if provided
        if "x_axis" in data:
            chart.category_axis.axis_title.text_frame.text = data["x_axis"]
        if "y_axis" in data:
            chart.value_axis.axis_title.text_frame.text = data["y_axis"]

        return chart